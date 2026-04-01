#!/usr/bin/env python3
"""Generate GaussCtrl thesis presentation slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY          = RGBColor(0x1F, 0x38, 0x64)
LIGHT_BLUE    = RGBColor(0xD6, 0xE4, 0xF0)
MID_BLUE      = RGBColor(0x2E, 0x75, 0xB6)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x1A, 0x1A, 0x1A)
RED_FILL      = RGBColor(0xFF, 0xD7, 0xD7)
RED_BORDER    = RGBColor(0xC0, 0x00, 0x00)
GREEN_FILL    = RGBColor(0xE2, 0xEF, 0xDA)
GREEN_BORDER  = RGBColor(0x37, 0x86, 0x30)
YELLOW_FILL   = RGBColor(0xFF, 0xF2, 0xCC)
YELLOW_BORDER = RGBColor(0xBF, 0x8F, 0x00)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
ORANGE_FILL   = RGBColor(0xFF, 0xE6, 0xCC)
ORANGE_BORDER = RGBColor(0xC5, 0x5A, 0x11)
SUBTLE_BLUE   = RGBColor(0xEA, 0xF4, 0xFB)


def make_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_bar(slide, prs, title, subtitle=None):
    h = Inches(1.4) if subtitle else Inches(1.1)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0), Inches(0), prs.slide_width, h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "   " + title
    run.font.color.rgb = WHITE
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = "   " + subtitle
        r2.font.color.rgb = RGBColor(0xBF, 0xD7, 0xED)
        r2.font.size = Pt(13)
        r2.font.name = "Calibri"
        r2.font.italic = True
    return h.inches


def add_box(slide, x, y, w, h, lines, fill=LIGHT_BLUE, border=MID_BLUE,
            font_size=12, bold_lines=None, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bold_lines = bold_lines or []
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT
        run.font.name = "Calibri"
        run.font.bold = (i in bold_lines)
    return box


def add_arrow_h(slide, x, y, w=0.5):
    """Horizontal arrow as a text box."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y - 0.2), Inches(w), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2192"
    run.font.size = Pt(24)
    run.font.color.rgb = MID_BLUE
    run.font.bold = True


def add_arrow_v(slide, cx, y, h=0.38):
    """Vertical arrow as a text box, cx is horizontal centre."""
    tb = slide.shapes.add_textbox(Inches(cx - 0.25), Inches(y), Inches(0.5), Inches(h))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "\u2193"
    run.font.size = Pt(20)
    run.font.color.rgb = MID_BLUE
    run.font.bold = True


def add_textbox(slide, x, y, w, h, text, font_size=11, bold=False,
                color=DARK_TEXT, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = "Calibri"
    run.font.bold = bold
    run.font.italic = italic


# ── Slide 1: Original GaussCtrl ───────────────────────────────────────────────
def slide1_original(prs):
    sld = blank_slide(prs)
    add_title_bar(sld, prs, "Original GaussCtrl: Multi-View Editing Flow")

    BOX_W, BOX_H = 2.7, 2.1
    ARR_W = 0.45
    TOTAL = 4 * BOX_W + 3 * ARR_W
    sx = (13.33 - TOTAL) / 2
    sy = 1.85

    data = [
        (["1. Render", "Training Views", "(RGB + Depth)"],             LIGHT_BLUE,   MID_BLUE),
        (["2. DDIM Inversion", "\u2192 noisy latent z\u2080", "per view"], LIGHT_BLUE,   MID_BLUE),
        (["3. Edit in Chunks", "ControlNet + Depth", "Cross-View Attention"], ORANGE_FILL, ORANGE_BORDER),
        (["4. Fine-tune 3DGS", "on edited images"],                    LIGHT_BLUE,   MID_BLUE),
    ]

    for i, (lines, fill, border) in enumerate(data):
        bx = sx + i * (BOX_W + ARR_W)
        add_box(sld, bx, sy, BOX_W, BOX_H, lines,
                fill=fill, border=border, font_size=13)
        if i < 3:
            add_arrow_h(sld, bx + BOX_W, sy + BOX_H / 2 - 0.02, w=ARR_W)

    # Cross-view attention note
    note_y = sy + BOX_H + 0.55
    add_box(sld, 0.6, note_y, 12.13, 1.4,
            ["Cross-View Attention:",
             "Each view attends to tokens from 4 fixed reference views "
             "(unedited, DDIM-inverted latents). Sharing attention to the same "
             "references enforces multi-view consistency throughout the diffusion process."],
            fill=LIGHT_GRAY, border=MID_BLUE, font_size=12,
            bold_lines=[0], align=PP_ALIGN.LEFT)


# ── Slide 2: First Attempt ─────────────────────────────────────────────────────
def slide2_first_attempt(prs, image_path):
    sld = blank_slide(prs)
    add_title_bar(sld, prs,
                  "First Attempt: Editing Reference View Without Cross-View Context",
                  subtitle="Branch: Sequential_Reference_Generation")

    bw, bh = 6.1, 1.05
    bx = 0.45
    top = 1.55

    add_box(sld, bx, top, bw, bh,
            ["Idea: edit ref\u2080 alone (no cross-view attention), "
             "DDIM-invert, use as reference latent for remaining views"],
            fill=LIGHT_BLUE, border=MID_BLUE, font_size=12, align=PP_ALIGN.LEFT)

    add_arrow_v(sld, bx + bw / 2, top + bh, h=0.38)

    add_box(sld, bx, top + bh + 0.38, bw, 2.5,
            ["Problem",
             "Without cross-view attention, ref\u2080 has no multi-view context.",
             "The model edits in isolation \u2192 appearance is inconsistent",
             "with the rest of the scene geometry and viewpoint.",
             "Result: unusable edit quality  (see image \u2192)"],
            fill=RED_FILL, border=RED_BORDER, font_size=12,
            bold_lines=[0], align=PP_ALIGN.LEFT)

    # Image
    ix, iy, iw, ih = 7.0, 1.55, 5.8, 5.3
    sld.shapes.add_picture(image_path, Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    add_textbox(sld, ix, iy + ih + 0.05, iw, 0.35,
                "ref\u2080 edited without cross-view attention  (experiment: panda_seq_v3)",
                font_size=10, italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ── Slide 3: Current Approach ──────────────────────────────────────────────────
def slide3_current(prs):
    sld = blank_slide(prs)
    add_title_bar(sld, prs,
                  "Current Approach: Sequential Reference Editing with Re-inversion",
                  subtitle="Branch: Sequential_Reference_z0")

    bw, bh = 10.0, 0.72
    gap = 0.33
    badge_col_w = 0.55
    total_w = badge_col_w + bw
    sx = (13.33 - total_w) / 2
    gx = sx + badge_col_w
    cur_y = 1.6

    steps = [
        ("1", "Edit ref\u2080 with cross-view attention to all 4 original (unedited) latents",
         LIGHT_BLUE, MID_BLUE, False),
        ("2", "DDIM-invert edited ref\u2080  \u2192  z\u2080 (edited latent)",
         SUBTLE_BLUE, MID_BLUE, False),
        ("3", "Edit ref\u2081 attending to the edited ref\u2080 latent",
         LIGHT_BLUE, MID_BLUE, False),
        ("4", "Continue: each ref_k attends to already-edited ref\u2080 \u2026 ref_{k\u22121}",
         LIGHT_BLUE, MID_BLUE, False),
        ("5", "Edit all target views attending to the 4 edited reference latents",
         GREEN_FILL, GREEN_BORDER, True),
    ]

    for i, (num, text, fill, border, is_final) in enumerate(steps):
        # Badge
        badge_x = sx
        badge_y = cur_y + (bh - 0.44) / 2
        badge = sld.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(badge_x), Inches(badge_y), Inches(0.44), Inches(0.44)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN_BORDER if is_final else NAVY
        badge.line.fill.background()
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = num
        br.font.color.rgb = WHITE
        br.font.size = Pt(12)
        br.font.bold = True
        br.font.name = "Calibri"

        add_box(sld, gx, cur_y, bw, bh, [text],
                fill=fill, border=border, font_size=12, align=PP_ALIGN.LEFT)

        if i < len(steps) - 1:
            add_arrow_v(sld, gx + bw / 2, cur_y + bh, h=gap)

        cur_y += bh + gap

    # Key insight
    insight_y = cur_y + 0.1
    add_box(sld, 0.5, insight_y, 12.33, 0.82,
            ["Key insight:  by propagating edits through reference views before editing "
             "target views, each target attends to already-edited (mutually consistent) "
             "references \u2014 improving cross-view coherence."],
            fill=YELLOW_FILL, border=YELLOW_BORDER, font_size=11, align=PP_ALIGN.LEFT)


# ── Main ───────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    IMAGE_PATH = ("/data/leuven/385/vsc38511/outputs/debug_edited_images/"
                  "panda_seq_v3/ref_00_idx0004_edited.png")
    OUTPUT = "/data/leuven/385/vsc38511/gaussctrl-fork/gaussctrl_presentation.pptx"

    prs = make_prs()
    slide1_original(prs)
    slide2_first_attempt(prs, IMAGE_PATH)
    slide3_current(prs)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
